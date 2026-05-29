from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation


def _shape_summary(shape) -> str:
    parts: list[str] = []
    parts.append(f"type={shape.shape_type}")
    if shape.has_text_frame:
        txt = (shape.text or "").strip().replace("\n", "\\n")
        if txt:
            parts.append(f"text='{txt[:80] + ('…' if len(txt) > 80 else '')}'")
    try:
        ph = shape.placeholder_format  # raises if not a placeholder
        parts.append(f"ph.type={ph.type}")
        parts.append(f"ph.idx={ph.idx}")
    except Exception:
        pass
    try:
        parts.append(f"name='{shape.name}'")
    except Exception:
        pass
    return ", ".join(parts)


def main() -> int:
    if len(sys.argv) != 2:
        print("Usage: ppt_inspect_template.py <template.pptx>")
        return 2

    template = Path(sys.argv[1])
    prs = Presentation(str(template))

    print(f"file: {template}")
    print(f"slides: {len(prs.slides)}")
    print(f"slide_size: {prs.slide_width} x {prs.slide_height}")
    print()

    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        try:
            if slide.shapes.title is not None:
                title = (slide.shapes.title.text or "").strip()
        except Exception:
            pass
        print(f"--- slide {i} ---")
        if title:
            print(f"title: {title}")
        print(f"shapes: {len(slide.shapes)}")
        for j, shape in enumerate(slide.shapes, start=1):
            print(f"{j:02d}. {_shape_summary(shape)}")
        print()

    return 0


if __name__ == "__main__":
    raise SystemExit(main())

