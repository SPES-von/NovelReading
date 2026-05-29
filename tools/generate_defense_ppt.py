from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation


def set_shape_text(shape, text: str) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False
    tf = shape.text_frame
    tf.clear()
    first = True
    for line in text.split("\n"):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = line
    return True


def replace_if_exact(shape, old: str, new: str) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False
    cur = (shape.text or "").strip()
    if cur == old:
        set_shape_text(shape, new)
        return True
    return False


def replace_if_contains(shape, needle: str, new: str) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False
    cur = (shape.text or "")
    if needle in cur:
        set_shape_text(shape, new)
        return True
    return False


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--template", required=True, help="Path to template .pptx")
    ap.add_argument("--output", required=True, help="Path to output .pptx")
    ap.add_argument("--title", default="基于 Vue 与 Spring Boot 的轻小说阅读系统设计与实现")
    ap.add_argument("--presenter", default="xxxx")
    ap.add_argument("--advisor", default="xxxx")
    ap.add_argument("--school", default="广州软件学院")
    args = ap.parse_args()

    prs = Presentation(args.template)

    # Slide contents (keep short: template says 2-3 pages for implementation)
    slide3 = (
        "系统分析（需求与用例）\n"
        "- 角色：访客、注册用户、系统管理员\n"
        "- 访客：浏览首页推荐、分类筛选、搜索、阅读章节、查看详情\n"
        "- 用户：登录/注册、书架管理、阅读进度同步、阅读偏好、听书（TTS）\n"
        "- 管理员：用户管理、小说上架/章节维护、分类/标签、评论审核、统计与设置"
    )
    slide4 = (
        "系统设计（架构/模块/数据库）\n"
        "- 架构：Vue 3 + Spring Boot 前后端分离，RESTful API（/api、/api/admin）\n"
        "- 安全：Spring Security + JWT，无状态鉴权；密码 BCrypt\n"
        "- 缓存：Redis 缓存热点数据与阅读进度（可选 no-redis 配置）\n"
        "- 数据库：MySQL 8，核心实体：小说/章节/作者/文库/标签/书架/进度/偏好"
    )
    slide5 = (
        "系统实现（核心功能）\n"
        "- 首页聚合：特色推荐、强推榜、新书榜、热点榜、完本推荐、热门小说等\n"
        "- 小说详情：封面/简介/标签/目录；加入书架、送鲜花、开始阅读\n"
        "- 阅读页：章节加载、上下章切换、目录跳转；主题/字体/字号等阅读设置\n"
        "- 认证：登录/注册签发 JWT；Axios 拦截器自动携带 Authorization"
    )
    slide6 = (
        "系统测试（功能测试）\n"
        "- 登录/注册：正确/错误账号密码、必填校验、401 提示\n"
        "- 首页与推荐：模块数据加载、跳转详情、交互验证\n"
        "- 阅读链路：章节加载、翻章、阅读设置保存与应用\n"
        "- 书架/个人中心：加入/移出书架、资料与偏好保存"
    )
    slide7 = (
        "系统测试（非功能测试）\n"
        "- 性能：首页聚合、列表筛选、章节接口多次请求，响应约 2–3 秒（手工观测）\n"
        "- 兼容性：Chrome/Edge/Firefox/360/搜狗/UC 等主要页面展示与操作正常\n"
        "- 安全：未登录访问受保护接口返回 401；后台接口权限隔离（403）"
    )
    slide8 = (
        "特色与创新\n"
        "- 前后端分离 + RESTful：界面与业务解耦，便于迭代与扩展\n"
        "- JWT 无状态鉴权 + Spring Security：接口权限清晰、跨端友好\n"
        "- Redis 缓存热点：降低数据库压力，提升首页/榜单等高频访问体验\n"
        "- 百度语音合成（TTS）：实现“看书 + 听书”场景，提升无障碍体验\n"
        "- 配置驱动榜单/推荐位：通过配置表维护展示顺序，运营维护成本低"
    )
    slide9 = (
        "不足与改进\n"
        "- 管理端部分模块（统计/系统设置等）当前以占位或待完善为主\n"
        "- 推荐以配置驱动为主，后续可引入基于标签/行为的个性化推荐\n"
        "- 听书可扩展：音色/语速、长文本分片、断点续播与缓存策略\n"
        "- 补充自动化与压力测试：JMeter/接口测试/端到端回归，提高质量与可靠性"
    )

    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            # Cover
            if idx == 1:
                replace_if_exact(shape, "《毕业论文题目》", f"《{args.title}》")
                replace_if_contains(shape, "答辩人：", f"答辩人：{args.presenter}")
                replace_if_contains(shape, "指导教师：", f"指导教师：{args.advisor}")
                replace_if_exact(shape, "广州软件学院", args.school)

            # Section titles / fixes
            if idx == 4:
                replace_if_exact(shape, "系统系统设计", "系统设计")
            if idx == 6:
                replace_if_exact(shape, "系统测试", "系统测试（功能）")
            if idx == 7:
                replace_if_exact(shape, "系统测试", "系统测试（非功能）")

            # Replace template "content requirements" boxes
            if idx == 3:
                replace_if_contains(shape, "（内容要求：", slide3)
            elif idx == 4:
                replace_if_contains(shape, "（内容要求：", slide4)
            elif idx == 5:
                replace_if_contains(shape, "（内容要求：", slide5)
            elif idx == 6:
                replace_if_contains(shape, "（内容要求：", slide6)
            elif idx == 7:
                replace_if_contains(shape, "（内容要求：", slide7)

            # Slide 8 and 9: large rectangle, no default text; inject into the main text area by name.
            if idx == 8 and getattr(shape, "has_text_frame", False):
                # Only fill the big rectangle (AutoShape), avoid overwriting the title
                if shape.name.startswith("矩形"):
                    set_shape_text(shape, slide8)
            if idx == 9 and getattr(shape, "has_text_frame", False):
                # Template has an empty text box named 文本框 7
                if shape.name == "文本框 7" or (shape.text or "").strip() == "":
                    set_shape_text(shape, slide9)

            # End slide personal info
            if idx == 11:
                replace_if_contains(shape, "答辩人：", f"答辩人：{args.presenter}")
                replace_if_contains(shape, "指导教师：", f"指导教师：{args.advisor}")
                replace_if_exact(shape, "广州软件学院", args.school)

    out = Path(args.output)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

