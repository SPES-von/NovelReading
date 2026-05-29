from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation


def shape_to_dict(shape) -> dict:
    d: dict = {
        "name": getattr(shape, "name", ""),
        "shape_type": int(getattr(shape, "shape_type", 0) or 0),
        "has_text_frame": bool(getattr(shape, "has_text_frame", False)),
        "is_placeholder": bool(getattr(shape, "is_placeholder", False)),
    }
    if d["has_text_frame"]:
        d["text"] = (shape.text or "").strip()
    if d["is_placeholder"]:
        try:
            d["placeholder_type"] = str(shape.placeholder_format.type)
        except Exception:
            d["placeholder_type"] = None
    return d


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--template", required=True, help="Path to template .pptx")
    ap.add_argument("--pretty", action="store_true", help="Pretty print JSON")
    args = ap.parse_args()

    template = Path(args.template)
    prs = Presentation(str(template))

    out = {
        "template": str(template),
        "slide_count": len(prs.slides),
        "slides": [],
    }

    for idx, slide in enumerate(prs.slides, start=1):
        slide_obj = {
            "index": idx,
            "shapes": [shape_to_dict(s) for s in slide.shapes],
        }
        out["slides"].append(slide_obj)

    if args.pretty:
        print(json.dumps(out, ensure_ascii=False, indent=2))
    else:
        print(json.dumps(out, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

